import os
import re
import unicodedata
import PyPDF2
from pptx import Presentation
from .config import ALLOWED_EXTENSIONS


def allowed_file(filename):
    """检查文件类型是否被允许"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def extract_text_from_file(filepath):
    """
    根据文件扩展名提取文本内容
    """

    
    # 确保filepath是字符串类型
    if not isinstance(filepath, str):
        raise ValueError(f"文件路径必须是字符串类型，当前类型: {type(filepath)}")
    
    # 检查文件是否存在
    if not os.path.exists(filepath):
        raise ValueError(f"文件不存在: {filepath}")
    
    # 获取文件名部分（去除路径）
    filename = os.path.basename(filepath)
    
    # 处理中文文件名编码问题
    try:
        # 确保文件名是正确的Unicode字符串
        if isinstance(filename, bytes):
            filename = filename.decode('utf-8')
        
        # 标准化文件名
        filename = unicodedata.normalize('NFC', filename)
        
    except Exception as e:
        print(f"文件名编码处理错误: {e}")
        # 如果编码处理失败，尝试使用原始文件名
        pass
    
    # 检查文件名是否包含扩展名
    if '.' not in filename:
        raise ValueError(f"文件没有扩展名: {filename}")
    
    # 使用更安全的方式提取文件扩展名
    try:
        # 分割文件名和扩展名
        name_parts = filename.rsplit('.', 1)
        if len(name_parts) != 2:
            raise ValueError(f"无法解析文件扩展名: {filename}")
        
        file_extension = name_parts[1].lower().strip()
        print(f"提取的扩展名: '{file_extension}'")
        
        # 验证扩展名不为空且有效
        if not file_extension or len(file_extension) == 0:
            raise ValueError(f"文件扩展名为空: {filename}")
        
        # 移除可能的特殊字符
        file_extension = re.sub(r'[^a-zA-Z0-9]', '', file_extension)
        print(f"清理后的扩展名: '{file_extension}'")
        
    except Exception as e:
        print(f"扩展名提取错误: {e}")
        raise ValueError(f"无法提取文件扩展名: {filename}")
    
    if file_extension in ['ppt', 'pptx']:
        return extract_text_from_ppt(filepath)
    elif file_extension == 'pdf':
        return extract_text_from_pdf(filepath)
    elif file_extension == 'txt':
        return extract_text_from_txt(filepath)
    else:
        raise ValueError(f"不支持的文件格式: {file_extension}")


def extract_text_from_ppt(filepath):
    """
    从PPT文件提取文本
    """
    try:
        prs = Presentation(filepath)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_runs.append(shape.text.strip())
        return "\n".join(text_runs)
    except Exception as e:
        raise Exception(f"PPT文件处理失败: {str(e)}")


def extract_text_from_pdf(filepath):
    """
    从PDF文件提取文本
    """
    try:
        text_content = []
        with open(filepath, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                if text.strip():
                    text_content.append(text.strip())
        return "\n".join(text_content)
    except Exception as e:
        raise Exception(f"PDF文件处理失败: {str(e)}")


def extract_text_from_txt(filepath):
    """
    从文本文件提取内容
    """
    try:
        # 尝试多种编码格式
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
        for encoding in encodings:
            try:
                with open(filepath, 'r', encoding=encoding) as file:
                    content = file.read().strip()
                    if content:
                        return content
            except UnicodeDecodeError:
                continue
        raise Exception("无法读取文本文件，可能是编码问题")
    except Exception as e:
        raise Exception(f"文本文件处理失败: {str(e)}")