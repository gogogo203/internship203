import os
import json
import re
import requests
from flask import Flask, render_template, request
from pptx import Presentation
from dotenv import load_dotenv

app = Flask(__name__)
load_dotenv()

API_URL = "https://router.huggingface.co/v1/chat/completions"
HEADERS = {"Authorization": f"Bearer {os.environ['HF_TOKEN']}"}
UPLOAD_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload():
    file = request.files.get('pptfile')
    if not file or file.filename == '':
        return "请上传PPT文件"

    filepath = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
    file.save(filepath)

    text = extract_text_from_ppt(filepath)
    questions = generate_mcq_from_text(text)

    # 给每道题加上选项标签对应关系，传给模板方便渲染
    for q in questions:
        q['options_with_labels'] = list(zip(['A','B','C','D'], q['options']))

    return render_template('quiz.html', questions=questions)

def extract_text_from_ppt(filepath):
    prs = Presentation(filepath)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def query(payload):
    response = requests.post(API_URL, headers=HEADERS, json=payload)
    return response.json()

def generate_mcq_from_text(text):
    if not text or not isinstance(text, str):
        return []

    prompt = f"""
你是一个专业教育助手，请根据下面的讲稿内容生成 5 道有挑战性的选择题。
每题必须有4个选项，选项描述中不要含有ABCD字样，且只有一个正确答案。
请用 JSON 格式返回，格式如下：
[
    {{
        "question": "问题文本",
        "options": ["选项A", "选项B", "选项C", "选项D"],
        "answer": "A"
    }},
    ...
]

讲稿内容：
{text}
"""
    payload = {
        "model": "deepseek-ai/DeepSeek-R1:novita",
        "messages": [
            {"role": "system", "content": "你是专业出题助手"},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.7,
        "max_tokens": 1000
    }

    resp = query(payload)
    raw_content = resp.get("choices", [{}])[0].get("message", {}).get("content", "")
    # 提取json部分
    match = re.search(r"```json(.*?)```", raw_content, re.DOTALL)
    if match:
        json_str = match.group(1).strip()
    else:
        json_str = raw_content.strip()

    try:
        quiz_list = json.loads(json_str)
    except Exception as e:
        print("JSON解析错误:", e)
        quiz_list = []
    for q in quiz_list:
        q['options_with_labels'] = list(zip(['A', 'B', 'C', 'D'], q['options']))
    return quiz_list

@app.route('/submit', methods=['POST'])
def submit():
    total = int(request.form.get('total', 0))
    correct = 0
    results = []
    for i in range(total):
        user_ans = request.form.get(f'q{i}')
        correct_ans = request.form.get(f'answer{i}')
        is_correct = user_ans == correct_ans
        if is_correct:
            correct += 1
        results.append({
            'index': i + 1,
            'user': user_ans,
            'correct': correct_ans,
            'is_correct': is_correct
        })
    score = f"{correct} / {total}"
    return render_template('quiz_result.html', results=results, score=score)

if __name__ == '__main__':
    app.run(debug=True)
